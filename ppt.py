from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Define a title and content slide layout
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

def add_slide(title, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    tf = content_placeholder.text_frame
    tf.clear()
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)

# Slide 1: Title Slide
slide_0 = prs.slides.add_slide(title_slide_layout)
slide_0.shapes.title.text = "MFPCL Website Improvement Proposal"
slide_0.placeholders[1].text = "Presented by Shyanukant Rathi"

# Slide 2: Welcome & Objective
add_slide("Welcome & Objective", [
    "✔ Review your current website & digital presence",
    "✔ Showcase a live demo of proposed redesign",
    "✔ Offer tailored solutions: Web Design, AI Chatbot, and Social Media Management",
    "✔ Discuss pricing, discounts & FREE bonuses"
])

# Slide: Key Website Issues
add_slide("⚠️ Key Website Issues Identified", [
    "1. Website Performance (Speed)",
    "   - Slow load times on desktop and mobile",
    "   - Causes frustration and higher bounce rates",
    "2. Accessibility Problems",
    "   - Missing ARIA labels and contrast issues",
    "   - Not compliant with WCAG accessibility standards",
    "3. Broken or Missing Links",
    "   - Internal navigation links not working",
    "   - No social media links present",
    "4. Missing FAQ Section",
    "   - No easy-to-find answers for common user queries",
    "5. No Chatbot Support",
    "   - No real-time help or user engagement tool available"
])

# Slide: Proposed Solutions
add_slide("🚀 Proposed Solutions", [
    "✅ Website Performance Optimization",
    "   - Compress and lazy-load images",
    "   - Minify JavaScript & CSS",
    "   - Enable browser caching",
    "✅ Accessibility Enhancements",
    "   - Add proper ARIA labels",
    "   - Improve color contrast",
    "   - Ensure keyboard navigation support",
    "✅ Fix Navigation & Social Links",
    "   - Repair broken internal links",
    "   - Integrate active social media accounts",
    "✅ Add an FAQ Section",
    "   - Design a clean, searchable FAQ area",
    "   - Use real user questions to boost SEO and usability",
    "✅ Implement a Smart Chatbot",
    "   - 24/7 support with instant replies",
    "   - Can answer FAQs, direct users to pages, and collect leads"
])

# Slide 3: Website Design Packages (More Detailed)
add_slide("🌐 Website Design & Development", [
    "🔹 Basic Website: ₹15,000 (New) / ₹10,000 (Redesign)",
    "    - Portfolio / Personal Branding Site",
    "    - Landing Page for Leads or Events",
    "    - Blog with Search & Tags",
    "    - Mobile Responsive + Contact Form",
    "    - Basic On-Page SEO + SSL Integration",
    "",
    "🔹 Advanced Website: ₹25,000 (New) / ₹20,000 (Redesign)",
    "    - All Basic Features +",
    "    - E-commerce Functionality (Product Pages, Cart, Payment Integration)",
    "    - Online Service Booking / Training Webinars",
    "    - LMS for Paid Courses, Certificates",
    "    - Email Automation (Leads, Reminders)",
    "    - Sales Funnel Integration (Landing → Checkout → Thank You)"
])

# Slide 4: AI Chatbot Integration
add_slide("🤖 AI Chatbot Integration", [
    "🔹 Basic Chatbot (FAQ): ₹5,000 (One-Time Setup)",
    "    - Answers to common questions",
    "    - Works 24/7 without human intervention",
    "",
    "🔹 Knowledge Base Chatbot: ₹10,000 Setup",
    "    - Integrated with your content and knowledge base",
    "    - Learns from your website, PDFs, documents",
    "    - Smarter responses, interactive experience",
    "    + ₹2,000/month for:",
    "        - GPT API usage",
    "        - Hosting and Server Maintenance",
    "        - Monthly Performance Monitoring"
])

# Slide 5: Social Media Management
add_slide("📱 Social Media Management", [
    "🔹 One-time Setup (FB & IG): ₹10,000",
    "    - Page/Profile Optimization",
    "    - Content Calendar for 1 Month",
    "    - Automated Scheduled Posts",
    "",
    "🔹 Monthly Social Media Management: ₹15,000",
    "    - Platforms: Instagram, Facebook",
    "    - Reels & Short Videos (4-6/month)",
    "    - Analytics Reports & Engagement Growth",
    "",
    "🔹 Premium Plan (₹20,000/month)",
    "    - Platforms: IG + FB + YouTube",
    "    - Includes Video Shorts + YouTube Uploads",
    "    - Content Planning, Hashtag Research, Editing"
])

# Slide 6: Bundle Discounts
add_slide("💸 Bundle Discounts", [
    "✔ Choose Any 2 Services – Get ₹3,000 OFF",
    "✔ Choose All 3 Services – Get ₹5,000 OFF",
    "✔ Extra ₹500 OFF if using Basic Website Package",
    "",
    "📊 Example:",
    "Basic Website + Basic Chatbot = ₹20,000 → ₹17,000",
    "Advanced Redesign + Chatbot + Social Media = ₹45,000 → ₹40,000"
])

# Save the presentation
pptx_path = "data/MFPCL_Proposal_Detailed_Services_Shyanukant.pptx"
prs.save(pptx_path)
pptx_path
